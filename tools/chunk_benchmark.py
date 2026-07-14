#!/usr/bin/env python
"""Chunking benchmark — compare 6 methods across structural + retrieval metrics.

Supports multi‑modal input: PDF, DOCX, PPTX, XLSX, Markdown, plain text, images (OCR).

Usage::

    python tools/chunk_benchmark.py --dir ./docs --query "问题1" "问题2"
    python tools/chunk_benchmark.py --files a.pdf b.docx --query "关键词"
    python tools/chunk_benchmark.py --files report.md            # structural only
"""

from __future__ import annotations

import argparse
import json
import logging
import re
import sys
import textwrap
import time
from dataclasses import dataclass, field
from functools import partial
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

# ── Multi‑format text extraction (standalone, no project imports needed) ──────

_SUPPORTED_SUFFIXES = {
    ".txt", ".md", ".markdown", ".csv", ".json",
    ".pdf", ".docx", ".pptx", ".ppt", ".xlsx", ".xls",
    ".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp",
}


def _extract_text(file_path: str) -> str:
    """Extract readable text from any supported file type."""
    p = Path(file_path)
    suf = p.suffix.lower()
    raw: str | None = None

    # ── PDF ──
    if suf == ".pdf":
        try:
            import fitz  # pymupdf
        except ImportError:
            pass
        else:
            parts = []
            try:
                doc = fitz.open(str(p))
                for page in doc:
                    parts.append(page.get_text())
                doc.close()
                raw = "\n\n".join(parts)
            except Exception:
                pass

    # ── DOCX ──
    elif suf == ".docx":
        try:
            from docx import Document
            doc = Document(str(p))
            raw = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            pass
        except Exception:
            pass

    # ── PPTX / PPT ─────────────────────────────────────────────────
    elif suf in {".pptx", ".ppt"}:
        if suf == ".ppt":
            raw = (
                f"[旧格式 PPT 不支持] {p.name}\n"
                "请用 PowerPoint 另存为 .pptx 格式后重试。"
            )
        else:
            try:
                from pptx import Presentation
                prs = Presentation(str(p))
                # First pass: collect all shape texts to detect duplicates
                all_texts: list[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            all_texts.append(shape.text_frame.text.strip())
                # Identify repeated elements (appear in > 30% of slides)
                from collections import Counter
                counts = Counter(all_texts)
                dup_threshold = max(3, len(prs.slides) * 0.3)
                repeated = {t for t, c in counts.items() if c > dup_threshold}

                parts: list[str] = []
                for si, slide in enumerate(prs.slides, start=1):
                    slide_lines = [f"## Slide {si}"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            txt = shape.text_frame.text.strip()
                            if txt and txt not in repeated:
                                slide_lines.append(txt)
                    if len(slide_lines) > 1:
                        parts.append("\n".join(slide_lines))
                raw = "\n\n".join(parts) if parts else (
                    "\n".join(t for t in all_texts if t not in repeated)
                    if all_texts else ""
                )
            except ImportError:
                raw = f"[缺少 python-pptx] {p.name}"
            except Exception:
                raw = f"[PPT 解析失败] {p.name}"

    # ── XLSX (multi‑sheet, each sheet becomes a markdown section) ──
    elif suf in {".xlsx", ".xls"}:
        try:
            import pandas as pd
            xls = pd.ExcelFile(str(p))
            parts: list[str] = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                cols = list(df.columns)
                preview_rows = min(len(df), 100)
                lines = [
                    f"## Sheet: {sheet_name}",
                    f"columns={cols}",
                    f"rows={len(df)}",
                    "",
                ]
                for _, row in df.head(preview_rows).iterrows():
                    line = " | ".join(
                        f"{cols[i] if i < len(cols) else '?'}={v}"
                        for i, v in enumerate(row.values)
                        if str(v).strip()
                    )
                    if line:
                        lines.append(line)
                parts.append("\n".join(lines))
            raw = "\n\n".join(parts)
        except ImportError:
            pass
        except Exception:
            pass

    # ── Image (RapidOCR) ──────────────────────────────────────────
    elif suf in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}:
        try:
            from rapidocr_onnxruntime import RapidOCR
            engine = RapidOCR()
            result, _ = engine(str(p))
            if result:
                lines = [line[1] for line in result if line[1]]
                raw = "\n".join(lines)
            else:
                raw = f"[图片无文字] {p.name}"
        except ImportError:
            raw = f"[图片文件] {p.name} (安装 RapidOCR 后可提取文字)"
        except Exception:
            raw = f"[图片OCR失败] {p.name}"

    # ── Plain text / Markdown ──
    if raw is None:
        try:
            raw = p.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            try:
                raw = p.read_bytes().decode("utf-8", errors="ignore")
            except Exception:
                raw = f"[无法读取] {p.name}"

    return raw.replace("\x00", "").strip()


# ═══════════════════════════════════════════════════════════════════════════════
# Chunking methods
# ═══════════════════════════════════════════════════════════════════════════════

def _norm(s: str) -> str:
    return re.sub(r"\s+", " ", s).strip()


# ── 1. 固定长度切分 ──────────────────────────────────────────────────────────

def chunk_fixed(text: str, chunk_size: int = 500) -> List[str]:
    """Slice text at fixed character boundaries (no overlap)."""
    if not text.strip():
        return []
    return [_norm(text[i:i + chunk_size]) for i in range(0, len(text), chunk_size)]


# ── 2. 滑动窗口切分 ──────────────────────────────────────────────────────────

def chunk_sliding(text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
    """Fixed‑size chunks with overlap between neighbours."""
    if not text.strip():
        return []
    step = max(1, chunk_size - overlap)
    chunks: List[str] = []
    for i in range(0, len(text) - chunk_size + 1, step):
        chunks.append(_norm(text[i:i + chunk_size]))
    # Tail
    if len(text) % step != 0:
        tail = text[-chunk_size:] if len(text) > chunk_size else text
        chunks.append(_norm(tail))
    return chunks or [_norm(text)]


# ── 3. 段落切分 ──────────────────────────────────────────────────────────────

def chunk_paragraph(text: str) -> List[str]:
    """Split at double‑newline (paragraph) boundaries."""
    if not text.strip():
        return []
    parts = re.split(r"\n\s*\n", text)
    return [_norm(p) for p in parts if p.strip()]


# ── 4. Recursive 切分 (LangChain 风格) ────────────────────────────────────────

_RECURSIVE_SEPARATORS = [
    "\n\n", "\n",
    r"(?<=[。！？])(?=\S)",
    r"(?<=[.!?])\s+(?=\S)",
    r"(?<=[；;])(?=\S)",
    r"(?<=[，,;])(?=\S)",
    " ",
]


def _recursive_split(text: str, separators: List[str], target: int) -> List[str]:
    """Recursively split *text* by *separators* until each piece ≤ *target*."""
    if len(text) <= target:
        return [text] if text.strip() else []

    for sep in separators:
        if isinstance(sep, str) and sep.startswith(r"(?<="):
            parts = re.split(sep, text)
        elif sep in text:
            parts = text.split(sep)
        else:
            continue

        good: List[str] = []
        for part in parts:
            part = part.strip()
            if not part:
                continue
            if len(part) <= target:
                good.append(part)
            else:
                good.extend(_recursive_split(part, separators[1:], target))
        if good:
            return good

    # Force‑split at target size
    return [text[i:i + target] for i in range(0, len(text), target)]


def chunk_recursive(text: str, chunk_size: int = 500, overlap: int = 0) -> List[str]:
    """Recursive character split with LangChain‑style separator hierarchy."""
    if not text.strip():
        return []
    chunks = _recursive_split(text, _RECURSIVE_SEPARATORS, chunk_size)
    if overlap > 0:
        overlapped = []
        for i, c in enumerate(chunks):
            if i > 0 and overlap > 0:
                prev_tail = chunks[i - 1]
                if len(prev_tail) > overlap:
                    prev_tail = prev_tail[-overlap:]
                c = prev_tail + c
            overlapped.append(_norm(c))
        return overlapped
    return [_norm(c) for c in chunks]


# ── 5. Header 切分 (Markdown 标题感知) ────────────────────────────────────────

_HEADER_RE = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)


def chunk_header(text: str) -> List[str]:
    """Split at Markdown header boundaries; each section becomes a chunk."""
    if not text.strip():
        return []
    # Find header positions
    positions: List[Tuple[int, str]] = []  # (start, header_line)
    for m in _HEADER_RE.finditer(text):
        positions.append((m.start(), m.group(0)))

    if not positions:
        return chunk_paragraph(text)  # fallback to paragraph split

    chunks: List[str] = []
    for i, (pos, header) in enumerate(positions):
        start = pos
        end = positions[i + 1][0] if i + 1 < len(positions) else len(text)
        chunks.append(_norm(text[start:end]))
    return chunks


# ── 6. 句子切分 (语义切分的简易替代) ──────────────────────────────────────────

_SENTENCE_RE = re.compile(r"[^。！？.!?\n]+[。！？.!?\n]?")


def chunk_sentence(text: str) -> List[str]:
    """Split at sentence boundaries (。！？.!?)."""
    if not text.strip():
        return []
    parts = _SENTENCE_RE.findall(text)
    chunks = [_norm(p) for p in parts if p.strip()]
    # Merge very short sentences (<20 chars) into the previous one
    merged: List[str] = []
    for c in chunks:
        if merged and len(c) < 20:
            merged[-1] = merged[-1] + c
        else:
            merged.append(c)
    return merged


# ═══════════════════════════════════════════════════════════════════════════════
# Benchmark engine
# ═══════════════════════════════════════════════════════════════════════════════

@dataclass
class MethodResult:
    name: str
    chunk_count: int = 0
    avg_length: float = 0.0
    min_length: int = 0
    max_length: int = 0
    para_cut_rate: float = 0.0   # 0‑1
    header_cross_rate: float = 0.0  # 0‑1
    recall_at_3: float = 0.0
    recall_at_5: float = 0.0
    mrr: float = 0.0
    retrieval_ms: float = 0.0
    embed_ms: float = 0.0
    chunks: List[str] = field(default_factory=list)



def _compute_para_cut_rate(chunks: List[str], full_text: str) -> float:
    """Fraction of chunks whose start/end falls *inside* a paragraph.

    Works on a whitespace‑normalised copy of *full_text* so chunk positions
    can be located reliably despite ``_norm()`` having been applied.
    """
    if not chunks:
        return 0.0

    norm_text = _norm(full_text)

    # Map paragraph-boundary positions from original → normalised text
    para_positions: List[int] = []
    for m in re.finditer(r"\n\s*\n", full_text):
        prefix = full_text[:m.start()]
        para_positions.append(len(_norm(prefix)))

    # Also add implied boundaries at very start / very end
    para_positions = sorted(set(para_positions))
    if not para_positions:
        return 0.0

    cut_count = 0
    search_from = 0
    for c in chunks:
        # Locate chunk in the normalised text
        pos = norm_text.find(c, search_from)
        if pos < 0:
            head = c[:60].strip()
            if head:
                pos = norm_text.find(head, search_from)
        if pos < 0:
            continue
        end_pos = pos + len(c)
        search_from = pos + 1

        # Determine which paragraph interval [pos, end_pos) falls into.
        # "Inside a paragraph" = chunk's start/end is not near any boundary.
        for pb in para_positions:
            if pb > pos + 5 and pb < end_pos - 5:
                cut_count += 1
                break

    return cut_count / len(chunks)


def _compute_header_cross_rate(chunks: List[str], full_text: str) -> float:
    """Proportion of adjacent chunk pairs that cross a Markdown header boundary.

    Returns a value in [0, 1] where *higher* = fewer crossings = better.
    """
    if len(chunks) < 2:
        return 0.0

    norm_text = _norm(full_text)

    # Map header positions from original → normalised text
    norm_header_positions: List[int] = []
    for m in _HEADER_RE.finditer(full_text):
        prefix = full_text[:m.start()]
        norm_header_positions.append(len(_norm(prefix)))

    if not norm_header_positions:
        return 0.0

    def _section_of(pos: int) -> int:
        for idx in range(len(norm_header_positions) - 1, -1, -1):
            if pos >= norm_header_positions[idx]:
                return idx
        return -1

    # Build a (start_pos) lookup for every chunk in norm_text
    chunk_starts: List[int] = []
    search_from = 0
    for c in chunks:
        pos = norm_text.find(c, search_from)
        if pos < 0:
            head = c[:60].strip()
            if head:
                pos = norm_text.find(head, search_from)
        if pos >= 0:
            chunk_starts.append(pos)
            search_from = pos + 1
        else:
            chunk_starts.append(-1)

    cross_count = 0
    valid_pairs = 0
    for i in range(len(chunk_starts) - 1):
        if chunk_starts[i] >= 0 and chunk_starts[i + 1] >= 0:
            valid_pairs += 1
            if _section_of(chunk_starts[i]) != _section_of(chunk_starts[i + 1]):
                cross_count += 1

    if valid_pairs == 0:
        return 0.0
    return 1.0 - (cross_count / valid_pairs)  # higher = better (fewer crosses)


def _embed_chunks_tfidf(chunks: List[str]):
    """TF‑IDF vectoriser as a zero‑dependency fallback."""
    try:
        from sklearn.feature_extraction.text import TfidfVectorizer
    except ImportError:
        raise RuntimeError("请安装 scikit-learn: pip install scikit-learn")
    vec = TfidfVectorizer(max_features=2000)
    return vec, vec.fit_transform(chunks)


def _cosine_sim(a, b):
    """Row‑wise cosine similarity between sparse/numpy matrices."""
    import numpy as np
    a_d = a.toarray() if hasattr(a, "toarray") else np.asarray(a)
    b_d = b.toarray() if hasattr(b, "toarray") else np.asarray(b)
    dot = np.dot(a_d, b_d.T)
    norm_a = np.linalg.norm(a_d, axis=1, keepdims=True) + 1e-10
    norm_b = np.linalg.norm(b_d, axis=1, keepdims=True).T + 1e-10
    return dot / (norm_a * norm_b)


def _keyword_relevance(query: str, chunks: List[str]) -> set:
    """Return indices of chunks that contain >= 50 % of the query's meaningful tokens.

    Tokens are extracted within *same‑script runs* (CJK / Latin / digit) so that
    "RAG技术" yields {"RA","AG"} ∪ {"技术"} rather than the useless "G技".
    """
    # ── Tokenize query into same‑script bigrams + unigrams ────────────
    import unicodedata

    def _script(c: str) -> str:
        """Classify a character into CJK / Latin / digit / other."""
        if "一" <= c <= "鿿" or "㐀" <= c <= "䶿":
            return "cjk"
        if c.isascii() and c.isalpha():
            return "latin"
        if c.isdigit():
            return "digit"
        return "other"

    tokens: set = set()
    i = 0
    q = query.strip()
    while i < len(q):
        sc = _script(q[i])
        j = i
        while j < len(q) and _script(q[j]) == sc:
            j += 1
        run = q[i:j]
        if sc in ("cjk", "latin") and len(run) >= 2:
            # Bigrams within this run
            for k in range(len(run) - 1):
                tokens.add(run[k:k + 2])
        elif run.strip():
            tokens.add(run)  # unigram fallback
        i = j

    if not tokens:
        return set()

    threshold = max(1, len(tokens) // 2)

    relevant: set = set()
    for idx, chunk in enumerate(chunks):
        hits = sum(1 for tok in tokens if tok in chunk)
        if hits >= threshold:
            relevant.add(idx)
    return relevant


def _run_retrieval(
    chunks: List[str],
    queries: List[str],
    vec,
    chunk_vecs,
) -> Tuple[float, float, float, float]:
    """Return (recall@3, recall@5, mrr, latency_ms).

    Pseudo‑relevance is defined by keyword (bigram) overlap — deterministic
    and method‑independent.  TF‑IDF ranking quality is then measured against
    this ground truth.
    """
    import numpy as np

    t0 = time.perf_counter()
    query_vecs = vec.transform(queries)
    sims = _cosine_sim(query_vecs, chunk_vecs)  # (queries × chunks)
    retrieval_ms = (time.perf_counter() - t0) * 1000

    if not chunks:
        return 0.0, 0.0, 0.0, retrieval_ms

    recalls_at_3: List[float] = []
    recalls_at_5: List[float] = []
    mrrs: List[float] = []

    for qi, (query, row) in enumerate(zip(queries, sims)):
        relevant = _keyword_relevance(query, chunks)
        if not relevant:
            continue

        ranked = np.argsort(-row)

        for k, store in [(3, recalls_at_3), (5, recalls_at_5)]:
            top_k = set(ranked[:k])
            hits = len(top_k & relevant)
            store.append(hits / min(k, len(relevant)))

        for rank, idx in enumerate(ranked, start=1):
            if idx in relevant:
                mrrs.append(1.0 / rank)
                break
        else:
            mrrs.append(0.0)

    recall3 = float(np.mean(recalls_at_3)) if recalls_at_3 else 0.0
    recall5 = float(np.mean(recalls_at_5)) if recalls_at_5 else 0.0
    mrr = float(np.mean(mrrs)) if mrrs else 0.0
    return recall3, recall5, mrr, retrieval_ms


# ═══════════════════════════════════════════════════════════════════════════════
# Main benchmark
# ═══════════════════════════════════════════════════════════════════════════════

CHUNK_SIZE = 500
OVERLAP = 100
TABLE_WIDTHS = [20, 8, 10, 8, 8, 10, 10, 10, 8, 8]


def _hdr(*cols: str) -> str:
    """Format a markdown table header."""
    header = "| " + " | ".join(f"{c:^{TABLE_WIDTHS[i]}}" for i, c in enumerate(cols)) + " |"
    sep = "|" + "|".join("-" * (w + 2) for w in TABLE_WIDTHS[:len(cols)]) + "|"
    return header + "\n" + sep


def _row(*cols: str) -> str:
    padded = []
    for i, c in enumerate(cols):
        w = TABLE_WIDTHS[i] if i < len(TABLE_WIDTHS) else 10
        padded.append(f"{c:^{w}}")
    return "| " + " | ".join(padded) + " |"


def _make_methods(
    chunk_size: int,
    overlap: int,
) -> List[Tuple[str, Callable[..., List[str]]]]:
    """Return (name, chunk_fn) pairs. Each fn takes a single text and returns chunks."""
    return [
        ("1.固定长度", partial(chunk_fixed, chunk_size=chunk_size)),
        ("2.滑动窗口", partial(chunk_sliding, chunk_size=chunk_size, overlap=overlap)),
        ("3.段落切分", chunk_paragraph),
        ("4.Recursive", partial(chunk_recursive, chunk_size=chunk_size, overlap=overlap)),
        ("5.Header",    chunk_header),
        ("6.句子切分",  chunk_sentence),
    ]


def run_benchmark(
    texts: Dict[str, str],
    queries: Optional[List[str]] = None,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = OVERLAP,
) -> List[MethodResult]:
    """Run all 6 chunking methods and return results.

    Each file is chunked independently; structural metrics are averaged across
    files; retrieval is evaluated per‑method on its own chunk corpus.
    """
    method_defs = _make_methods(chunk_size, overlap)
    results: List[MethodResult] = []

    for name, chunk_fn in method_defs:
        # ── Per‑file chunking ─────────────────────────────────────────
        per_file: List[Tuple[str, List[str]]] = []  # [(fname, chunks)]
        all_chunks: List[str] = []
        for fname, text in texts.items():
            chunks = chunk_fn(text)
            if chunks:
                per_file.append((fname, chunks))
                all_chunks.extend(chunks)

        r = MethodResult(name=name, chunks=all_chunks)
        r.chunk_count = len(all_chunks)
        if all_chunks:
            lengths = [len(c) for c in all_chunks]
            r.avg_length = sum(lengths) / len(all_chunks)
            r.min_length = min(lengths)
            r.max_length = max(lengths)

        # ── Structural metrics (averaged across files) ────────────────
        para_rates: List[float] = []
        header_rates: List[float] = []
        for fname, fchunks in per_file:
            orig = texts[fname]
            para_rates.append(_compute_para_cut_rate(fchunks, orig))
            header_rates.append(_compute_header_cross_rate(fchunks, orig))
        r.para_cut_rate = sum(para_rates) / len(para_rates) if para_rates else 0.0
        r.header_cross_rate = sum(header_rates) / len(header_rates) if header_rates else 0.0

        results.append(r)
        print(f"  ✓ {name}: {r.chunk_count} chunks, avg={r.avg_length:.0f} chars"
              f" (来自 {len(per_file)} 个文件)")

    # ── Retrieval metrics (per‑method, relative threshold) ────────────
    if queries and results:
        print("\n  向量化 + 检索中...")
        for r in results:
            if not r.chunks:
                continue
            t0 = time.perf_counter()
            try:
                vec, chunk_vecs = _embed_chunks_tfidf(r.chunks)
            except RuntimeError as exc:
                print(f"  ⚠ TF‑IDF 不可用: {exc}")
                return results
            r.embed_ms = (time.perf_counter() - t0) * 1000

            r3, r5, mrr, lat = _run_retrieval(r.chunks, queries, vec, chunk_vecs)
            r.recall_at_3 = r3
            r.recall_at_5 = r5
            r.mrr = mrr
            r.retrieval_ms = lat
            print(f"  ✓ {r.name}: R@3={r3:.3f} R@5={r5:.3f} MRR={mrr:.3f}")

    return results


def print_table(results: List[MethodResult]) -> None:
    """Print a Markdown comparison table."""
    cols = ["方法", "块数", "平均长度", "切断率%", "越头率%",
            "R@3", "R@5", "MRR", "检索ms", "嵌入ms"]
    print()
    print(_hdr(*cols))
    for r in results:
        print(_row(
            r.name,
            str(r.chunk_count),
            f"{r.avg_length:.0f}",
            f"{r.para_cut_rate * 100:.1f}",
            f"{(1 - r.header_cross_rate) * 100:.1f}" if r.header_cross_rate > 0 else "N/A",
            f"{r.recall_at_3:.3f}",
            f"{r.recall_at_5:.3f}",
            f"{r.mrr:.3f}",
            f"{r.retrieval_ms:.0f}" if r.retrieval_ms > 0 else "-",
            f"{r.embed_ms:.0f}" if r.embed_ms > 0 else "-",
        ))

    print()
    print("指标说明：")
    print("  切断率 = chunk 内部跨越段落边界的比例（越低越好）")
    print("  越头率 = 相邻 chunk 跨过 Markdown 标题边界的比例（越低越好）")
    print("  R@K   = Top-K 中命中相关 chunk 的比例（越高越好）")
    print("  MRR   = 第一个相关 chunk 排名的倒数均值（越高越好）")
    print()
    print("⚠ 相关性基于关键词命中（bigram 重叠≥50%），衡量 TF‑IDF 排序与关键词匹配的一致性。")


# ═══════════════════════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════════════════════

def main() -> None:
    # Ensure UTF‑8 output on Windows (PowerShell / cmd)
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    ap = argparse.ArgumentParser(
        description="Chunking benchmark — 6 methods × structural + retrieval metrics",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""\
        examples:
          python tools/chunk_benchmark.py --dir ./docs
          python tools/chunk_benchmark.py --files a.pdf b.docx --query "什么是RAG" "如何优化检索"
          python tools/chunk_benchmark.py --files report.md --chunk-size 800 --overlap 150
        """),
    )
    ap.add_argument("--dir", help="Directory containing documents to process")
    ap.add_argument("--files", nargs="+", default=[], help="Specific files to process")
    ap.add_argument("--query", nargs="+", default=[], help="Test queries for retrieval evaluation")
    ap.add_argument("--chunk-size", type=int, default=CHUNK_SIZE,
                    help=f"Target chunk size in characters (default: {CHUNK_SIZE})")
    ap.add_argument("--overlap", type=int, default=OVERLAP,
                    help=f"Overlap in characters for sliding/recursive (default: {OVERLAP})")
    ap.add_argument("--no-retrieval", action="store_true",
                    help="Skip retrieval evaluation (structural metrics only)")
    ap.add_argument("--verbose", "-v", action="store_true")
    args = ap.parse_args()

    if args.verbose:
        logging.basicConfig(level=logging.INFO)

    # ── Gather files ────────────────────────────────────────────────
    file_paths: List[Path] = [Path(f) for f in args.files]
    if args.dir:
        dir_path = Path(args.dir)
        if not dir_path.is_dir():
            print(f"错误: 目录不存在 — {args.dir}")
            sys.exit(1)
        file_paths.extend(sorted(dir_path.rglob("*")))

    file_paths = [p for p in file_paths if p.is_file() and p.suffix.lower() in _SUPPORTED_SUFFIXES]
    file_paths = list(dict.fromkeys(file_paths))  # dedup preserve order

    if not file_paths:
        print("错误: 未找到可处理的文件。支持的格式:", ", ".join(sorted(_SUPPORTED_SUFFIXES)))
        sys.exit(1)

    # ── Extract text ────────────────────────────────────────────────
    texts: Dict[str, str] = {}
    total_chars = 0
    print(f"正在从 {len(file_paths)} 个文件提取文本...")
    for fp in file_paths:
        content = _extract_text(str(fp))
        if content:
            texts[fp.name] = content
            total_chars += len(content)
            print(f"  ✓ {fp.name} ({len(content)} chars)")
        else:
            print(f"  ✗ {fp.name} (无内容)")

    if not texts:
        print("错误: 所有文件均无可提取的文本。")
        sys.exit(1)

    print(f"\n总计: {len(texts)} 个文件, {total_chars} 字符\n")

    # ── Run benchmark ───────────────────────────────────────────────
    queries = args.query if args.query and not args.no_retrieval else None
    if queries:
        print(f"测试查询 ({len(queries)}):")
        for q in queries:
            print(f"  · {q}")
        print()

    print("切分中...")
    results = run_benchmark(texts, queries, args.chunk_size, args.overlap)

    # ── Print results ───────────────────────────────────────────────
    print_table(results)

    # ── Also output JSON for programmatic use ───────────────────────
    json_path = Path("result") / "chunk_benchmark.json"
    json_path.parent.mkdir(parents=True, exist_ok=True)
    payload = []
    for r in results:
        d = {
            "method": r.name,
            "chunk_count": r.chunk_count,
            "avg_length": round(r.avg_length, 1),
            "min_length": r.min_length,
            "max_length": r.max_length,
            "para_cut_rate": round(r.para_cut_rate, 4),
            "header_cross_rate": round(r.header_cross_rate, 4),
            "recall_at_3": round(r.recall_at_3, 4),
            "recall_at_5": round(r.recall_at_5, 4),
            "mrr": round(r.mrr, 4),
            "retrieval_ms": round(r.retrieval_ms, 1),
            "embed_ms": round(r.embed_ms, 1),
        }
        payload.append(d)
    json_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"\n结果已保存至: {json_path}")


if __name__ == "__main__":
    main()
