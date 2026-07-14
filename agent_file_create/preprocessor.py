import logging
import os
import re
from difflib import SequenceMatcher
from io import BytesIO
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)


class _NullWriter:
    """File-like object that discards all output."""

    def write(self, s):
        pass

    def flush(self):
        pass


def _quiet_fitz():
    """Globally suppress PyMuPDF's stdout/stderr noise.
    find_tables prints "find_tables: exception occurred: ..." through
    pymupdf.message() which writes to various internal channels depending
    on the PyMuPDF version.  We try all known channels."""
    try:
        import fitz
        import pymupdf
        null = _NullWriter()
        # Try all known output channels across PyMuPDF versions
        for attr in ("_g_out_message", "_g_err_message", "_message_writer"):
            for mod in (fitz, pymupdf):
                try:
                    setattr(mod, attr, null)
                except Exception:
                    pass
    except Exception:
        pass


def _feature_enabled(name: str, default: bool = True) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return str(raw).strip().lower() in {"1", "true", "yes", "y", "on"}

# ---------------------------------------------------------------------------
# OCR
# ---------------------------------------------------------------------------

_OCR_INSTANCE: Optional[Any] = None
_EASYOCR_READER: Optional[Any] = None
_EASYOCR_ALLOWLIST = "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789 .,:;/-()$#@&%+*'\""


def get_ocr(use_cuda: bool = True) -> Any:
    """Return a cached RapidOCR instance (paddle GPU preferred, onnxruntime fallback)."""
    global _OCR_INSTANCE
    if _OCR_INSTANCE is not None:
        return _OCR_INSTANCE
    try:
        from rapidocr_paddle import RapidOCR
        _OCR_INSTANCE = RapidOCR(det_use_cuda=use_cuda, cls_use_cuda=use_cuda, rec_use_cuda=use_cuda)
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR
        _OCR_INSTANCE = RapidOCR()
    return _OCR_INSTANCE


def get_easyocr_reader() -> Any:
    """Return a cached EasyOCR reader when available."""
    global _EASYOCR_READER
    if _EASYOCR_READER is not None:
        return _EASYOCR_READER
    try:
        import easyocr
    except Exception:
        return None
    try:
        _EASYOCR_READER = easyocr.Reader(["en"], gpu=False)
    except Exception as e:
        logger.warning(f"easyocr_init_failed err={str(e)[:160]}")
        return None
    return _EASYOCR_READER


def ocr_image(image_data: bytes) -> str:
    """Run OCR on raw image bytes, return joined text lines."""
    try:
        ocr = get_ocr()
        from PIL import Image

        img = Image.open(BytesIO(image_data))
        # RapidOCR accepts image path, numpy array, or PIL Image
        result, _ = ocr(img)
        if not result:
            return ""
        lines = []
        for item in result:
            text = (item[1] or "").strip()
            if text:
                lines.append(text)
        return "\n".join(lines)
    except Exception as e:
        logger.warning(f"ocr_failed err={str(e)[:160]}")
        return ""


def ocr_image_with_boxes(image_data: bytes) -> list[dict[str, Any]]:
    """Run OCR and return structured line items with polygon boxes when available."""
    try:
        ocr = get_ocr()
        from PIL import Image

        img = Image.open(BytesIO(image_data))
        result, _ = ocr(img)
        if not result:
            return []

        items: list[dict[str, Any]] = []
        for entry in result:
            if not entry or len(entry) < 2:
                continue
            box = entry[0]
            text = str(entry[1] or "").strip()
            score = None
            if len(entry) >= 3:
                try:
                    score = float(entry[2])
                except Exception:
                    score = None
            if not text:
                continue

            xs: list[float] = []
            ys: list[float] = []
            if isinstance(box, (list, tuple)):
                for pt in box:
                    if isinstance(pt, (list, tuple)) and len(pt) >= 2:
                        try:
                            xs.append(float(pt[0]))
                            ys.append(float(pt[1]))
                        except Exception:
                            continue
            item = {
                "text": text,
                "box": box,
                "score": score,
            }
            if xs and ys:
                item["left"] = min(xs)
                item["right"] = max(xs)
                item["top"] = min(ys)
                item["bottom"] = max(ys)
                item["width"] = max(xs) - min(xs)
                item["height"] = max(ys) - min(ys)
                item["cx"] = (min(xs) + max(xs)) / 2.0
                item["cy"] = (min(ys) + max(ys)) / 2.0
            items.append(item)
        return items
    except Exception as e:
        logger.warning(f"ocr_with_boxes_failed err={str(e)[:160]}")
        return []


def easyocr_image_with_boxes(image_data: bytes) -> list[dict[str, Any]]:
    """Run EasyOCR and return structured line items with polygon boxes when available."""
    try:
        reader = get_easyocr_reader()
        if reader is None:
            return []
        from PIL import Image
        import numpy as np

        img = Image.open(BytesIO(image_data)).convert("RGB")
        result = reader.readtext(
            np.array(img),
            detail=1,
            paragraph=False,
            allowlist=_EASYOCR_ALLOWLIST,
        )
        items: list[dict[str, Any]] = []
        for entry in result or []:
            if not entry or len(entry) < 2:
                continue
            box = entry[0]
            text = str(entry[1] or "").strip()
            score = None
            if len(entry) >= 3:
                try:
                    score = float(entry[2])
                except Exception:
                    score = None
            if not text:
                continue
            xs: list[float] = []
            ys: list[float] = []
            if isinstance(box, (list, tuple)):
                for pt in box:
                    if isinstance(pt, (list, tuple)) and len(pt) >= 2:
                        try:
                            xs.append(float(pt[0]))
                            ys.append(float(pt[1]))
                        except Exception:
                            continue
            item = {"text": text, "box": box, "score": score, "engine": "easyocr"}
            if xs and ys:
                item["left"] = min(xs)
                item["right"] = max(xs)
                item["top"] = min(ys)
                item["bottom"] = max(ys)
                item["width"] = max(xs) - min(xs)
                item["height"] = max(ys) - min(ys)
                item["cx"] = (min(xs) + max(xs)) / 2.0
                item["cy"] = (min(ys) + max(ys)) / 2.0
            items.append(item)
        return items
    except Exception as e:
        logger.warning(f"easyocr_with_boxes_failed err={str(e)[:160]}")
        return []


def easyocr_image(image_data: bytes) -> str:
    items = easyocr_image_with_boxes(image_data)
    lines = [str(item.get("text") or "").strip() for item in items if str(item.get("text") or "").strip()]
    return "\n".join(lines)


def score_ocr_text_quality(text: str) -> float:
    """Heuristic OCR quality score for choosing between multiple variants."""
    s = str(text or "").strip()
    if not s:
        return 0.0

    lines = [line.strip() for line in s.splitlines() if line.strip()]
    words = re.findall(r"[A-Za-z]{2,}|\d+|[\u4e00-\u9fff]+", s)
    kv_hits = len(
        re.findall(
            r"(?i)\b(to|from|date|fax|phone|number|pages|page|address|note|special|instructions|sender)\b|[:：]",
            s,
        )
    )
    suspicious = len(re.findall(r"[^\w\s:：,./()\-&]", s))
    long_joined = sum(1 for line in lines if len(line) >= 18 and " " not in line and ":" not in line and "：" not in line)

    score = 0.0
    score += min(len(lines), 40) * 0.8
    score += min(len(words), 120) * 0.25
    score += min(kv_hits, 30) * 1.2
    score -= suspicious * 0.15
    score -= long_joined * 1.6
    return score


def choose_better_ocr_text(*texts: str) -> str:
    candidates = [str(t or "").strip() for t in texts if str(t or "").strip()]
    if not candidates:
        return ""
    best = candidates[0]
    best_score = score_ocr_text_quality(best)
    for candidate in candidates[1:]:
        score = score_ocr_text_quality(candidate)
        if score > best_score:
            best = candidate
            best_score = score
    return best


def merge_ocr_texts(primary: str, secondary: str, max_lines: int = 80) -> str:
    """Keep the stronger OCR text and append non-duplicate high-value lines from the backup."""
    first = str(primary or "").strip()
    second = str(secondary or "").strip()
    if not first:
        return second
    if not second:
        return first

    merged: list[str] = []
    seen: set[str] = set()
    for source in (first.splitlines(), second.splitlines()):
        for raw in source:
            line = raw.strip()
            if len(line) < 2:
                continue
            key = re.sub(r"\s+", " ", line).strip().lower()
            if key in seen:
                continue
            seen.add(key)
            merged.append(line)
            if len(merged) >= max_lines:
                return "\n".join(merged)
    return "\n".join(merged)


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------


def preprocess_text(text: str, max_chars: int = 8000) -> str:
    s = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    s = re.sub(r"\n{3,}", "\n\n", s)
    s = s.strip()
    if len(s) > max_chars:
        return s[:max_chars] + "…"
    return s


def read_text_file(file_path: str, max_chars: int = 12000) -> str:
    p = Path(file_path)
    data = p.read_bytes()
    try:
        return preprocess_text(data.decode("utf-8"), max_chars=max_chars)
    except UnicodeDecodeError:
        pass
    for enc in ("gb18030", "latin-1"):
        try:
            return preprocess_text(data.decode(enc), max_chars=max_chars)
        except Exception:
            continue
    return preprocess_text(data.decode("utf-8", errors="ignore"), max_chars=max_chars)


# ---------------------------------------------------------------------------
# Image preprocessing
# ---------------------------------------------------------------------------


def _cv_to_pil_gray(gray) -> Any:
    from PIL import Image

    return Image.fromarray(gray)


def _order_quad_points(points) -> Any:
    import numpy as np

    pts = np.array(points, dtype="float32")
    s = pts.sum(axis=1)
    diff = np.diff(pts, axis=1)
    ordered = np.zeros((4, 2), dtype="float32")
    ordered[0] = pts[np.argmin(s)]
    ordered[2] = pts[np.argmax(s)]
    ordered[1] = pts[np.argmin(diff)]
    ordered[3] = pts[np.argmax(diff)]
    return ordered


def _four_point_transform(image, points):
    import cv2
    import numpy as np

    rect = _order_quad_points(points)
    tl, tr, br, bl = rect
    width_a = np.linalg.norm(br - bl)
    width_b = np.linalg.norm(tr - tl)
    height_a = np.linalg.norm(tr - br)
    height_b = np.linalg.norm(tl - bl)
    max_width = max(1, int(max(width_a, width_b)))
    max_height = max(1, int(max(height_a, height_b)))

    dst = np.array(
        [
            [0, 0],
            [max_width - 1, 0],
            [max_width - 1, max_height - 1],
            [0, max_height - 1],
        ],
        dtype="float32",
    )
    matrix = cv2.getPerspectiveTransform(rect, dst)
    warped = cv2.warpPerspective(image, matrix, (max_width, max_height))
    return warped


def _detect_document_boundary(image: Any, min_area_ratio: float = 0.35) -> tuple[Any, bool]:
    """Detect a page-like quadrilateral and perspective-correct it when possible."""
    try:
        import cv2
        import numpy as np
    except Exception:
        return image, False

    rgb = np.array(image.convert("RGB"))
    gray = cv2.cvtColor(rgb, cv2.COLOR_RGB2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 60, 180)
    contours, _ = cv2.findContours(edges, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    if not contours:
        return image, False

    image_area = float(gray.shape[0] * gray.shape[1] or 1)
    best_quad = None
    best_area = 0.0
    for contour in sorted(contours, key=cv2.contourArea, reverse=True)[:15]:
        peri = cv2.arcLength(contour, True)
        approx = cv2.approxPolyDP(contour, 0.02 * peri, True)
        if len(approx) != 4:
            continue
        area = float(cv2.contourArea(approx))
        if area < image_area * min_area_ratio or area <= best_area:
            continue
        best_quad = approx.reshape(4, 2)
        best_area = area

    if best_quad is None:
        return image, False

    warped = _four_point_transform(rgb, best_quad)
    if warped is None or warped.size == 0:
        return image, False
    from PIL import Image

    return Image.fromarray(warped), True


def _deskew_image(image: Any, max_angle: float = 15.0) -> Any:
    """Detect and correct document skew using Hough line detection."""
    try:
        import cv2
        import numpy as np
    except Exception:
        return image

    gray = np.array(image.convert("L"))
    edges = cv2.Canny(gray, 50, 150, apertureSize=3)
    lines = cv2.HoughLinesP(
        edges,
        1,
        np.pi / 180,
        threshold=100,
        minLineLength=max(40, min(gray.shape[:2]) // 4),
        maxLineGap=20,
    )
    if lines is None:
        return image

    angles: list[float] = []
    for line in lines:
        x1, y1, x2, y2 = line[0]
        angle = float(np.degrees(np.arctan2(y2 - y1, x2 - x1)))
        while angle <= -90:
            angle += 180
        while angle > 90:
            angle -= 180
        if angle > 45:
            angle -= 90
        elif angle < -45:
            angle += 90
        if abs(angle) <= max_angle:
            angles.append(angle)

    if not angles:
        return image

    median_angle = float(np.median(angles))
    if abs(median_angle) < 0.5:
        return image

    h, w = gray.shape
    center = (w // 2, h // 2)
    matrix = cv2.getRotationMatrix2D(center, median_angle, 1.0)
    rotated = cv2.warpAffine(gray, matrix, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)
    return _cv_to_pil_gray(rotated)


def _adaptive_binarize(image: Any, block_size: int = 21, c: int = 12) -> Any:
    """Apply adaptive Gaussian thresholding for uneven lighting."""
    try:
        import cv2
        import numpy as np
    except Exception:
        return image

    gray = np.array(image.convert("L"))
    if block_size % 2 == 0:
        block_size += 1
    block_size = max(3, block_size)
    binary = cv2.adaptiveThreshold(
        gray,
        255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        block_size,
        c,
    )
    return _cv_to_pil_gray(binary)


def _denoise_nl_means(image: Any, h: float = 10.0) -> Any:
    """Remove noise while preserving text strokes."""
    try:
        import cv2
        import numpy as np
    except Exception:
        return image

    gray = np.array(image.convert("L"))
    denoised = cv2.fastNlMeansDenoising(gray, None, h, 7, 21)
    return _cv_to_pil_gray(denoised)


def _has_dense_texture(
    gray: Any,
    block_size: int = 32,
    threshold: float = 0.12,
) -> bool:
    """Return True when the image contains dense text-like regions.

    Scanned documents typically show moderate local variance in many blocks
    (text strokes), unlike photographs or empty pages.  Used to decide
    whether stronger binarization + denoising should be applied.
    """
    try:
        import cv2
        import numpy as np
    except Exception:
        return False

    arr = np.array(gray)
    h, w = arr.shape
    n_y = max(1, h // block_size)
    n_x = max(1, w // block_size)
    text_blocks = 0
    total = n_y * n_x

    for by in range(n_y):
        y1 = by * block_size
        y2 = min((by + 1) * block_size, h)
        for bx in range(n_x):
            x1 = bx * block_size
            x2 = min((bx + 1) * block_size, w)
            local_std = float(np.std(arr[y1:y2, x1:x2].astype(np.float64)))
            if 8.0 < local_std < 80.0:
                text_blocks += 1

    return (text_blocks / float(total or 1)) >= threshold


def preprocess_image_path(
    file_path: str,
    max_long_edge: int = 2048,
    jpeg_quality: int = 85,
    profile: str = "adaptive",
) -> bytes:
    try:
        from PIL import Image, ImageEnhance, ImageOps, ImageStat
    except Exception:
        return Path(file_path).read_bytes()

    try:
        img = Image.open(file_path)
        img = ImageOps.exif_transpose(img)
        img = img.convert("RGB")

        w, h = img.size
        long_edge = max(w, h)
        if long_edge > max_long_edge:
            ratio = max_long_edge / long_edge
            new_size = (max(1, int(w * ratio)), max(1, int(h * ratio)))
            img = img.resize(new_size, Image.LANCZOS)

        gray = ImageOps.grayscale(img)
        stat = ImageStat.Stat(gray)
        mean = float(stat.mean[0]) if stat.mean else 128.0
        std = float(stat.stddev[0]) if stat.stddev else 0.0

        profile_norm = str(profile or "adaptive").strip().lower()
        enable_boundary = _feature_enabled("Q1_ENABLE_BOUNDARY", True)
        enable_skew = _feature_enabled("Q1_ENABLE_SKEW", True)
        enable_denoise = _feature_enabled("Q1_ENABLE_DENOISE", True)
        enable_binarize = _feature_enabled("Q1_ENABLE_BINARIZE", True)

        boundary_corrected = False
        if enable_boundary and profile_norm == "adaptive" and std < 20:
            img, boundary_corrected = _detect_document_boundary(img)
        elif enable_boundary and profile_norm == "strong":
            img, boundary_corrected = _detect_document_boundary(img, min_area_ratio=0.45)

        gray = ImageOps.grayscale(img)
        stat = ImageStat.Stat(gray)
        mean = float(stat.mean[0]) if stat.mean else 128.0
        std = float(stat.stddev[0]) if stat.stddev else 0.0

        if profile_norm in {"light", "medium", "strong"}:
            mode = profile_norm
        else:
            if boundary_corrected or std < 20:
                mode = "photo"
            elif std < 50:
                mode = "strong"
            elif std < 65:
                # Check if image looks like a scanned document: if text-rich
                # regions are present, upgrade to strong for binarization.
                if enable_binarize and _has_dense_texture(gray):
                    mode = "strong"
                else:
                    mode = "medium"
            else:
                mode = "light"

        work = gray
        if mean < 85 or mean > 190:
            work = ImageOps.autocontrast(work, cutoff=2)

        if mode == "light":
            work = ImageOps.autocontrast(work, cutoff=1)
            work = ImageEnhance.Sharpness(work).enhance(1.15)
        elif mode == "medium":
            if enable_skew:
                work = _deskew_image(work)
            work = ImageOps.autocontrast(work, cutoff=2)
            work = ImageEnhance.Contrast(work).enhance(1.2)
            work = ImageEnhance.Sharpness(work).enhance(1.25)
        elif mode == "strong":
            if enable_skew:
                work = _deskew_image(work)
            if enable_denoise:
                work = _denoise_nl_means(work, h=12.0)
            if enable_binarize:
                work = _adaptive_binarize(work, block_size=21, c=12)
            work = ImageOps.autocontrast(work, cutoff=1)
        else:
            if enable_skew:
                work = _deskew_image(work)
            if enable_denoise:
                work = _denoise_nl_means(work, h=14.0)
            if enable_binarize:
                work = _adaptive_binarize(work, block_size=25, c=10)
            work = ImageOps.autocontrast(work, cutoff=1)

        buf = BytesIO()
        work.convert("RGB").save(buf, format="JPEG", quality=jpeg_quality)
        return buf.getvalue()
    except Exception as e:
        logger.warning(f"preprocess_image_failed err={str(e)[:160]}")
        return Path(file_path).read_bytes()


# ---------------------------------------------------------------------------
# PDF (PyMuPDF)
# ---------------------------------------------------------------------------


def extract_pdf_text_fast(file_path: str, max_chars: int = 20000) -> str:
    """Extract text layer from PDF using PyMuPDF (fitz)."""
    try:
        import fitz
        _quiet_fitz()
    except Exception:
        return _extract_pdf_text_pypdf2_fallback(file_path, max_chars)

    try:
        doc = fitz.open(str(file_path))
        parts: list[str] = []
        total = 0
        for page in doc:
            try:
                t = page.get_text("text") or ""
            except Exception:
                t = ""
            if t:
                parts.append(t)
                total += len(t)
                if total >= max_chars:
                    break
        doc.close()
        return preprocess_text("\n".join(parts), max_chars=max_chars)
    except Exception as e:
        logger.warning(f"pdf_text_extract_fitz_failed err={str(e)[:160]}")
        return _extract_pdf_text_pypdf2_fallback(file_path, max_chars)


def _extract_pdf_text_pypdf2_fallback(file_path: str, max_chars: int = 20000) -> str:
    """Fallback: extract text using PyPDF2."""
    try:
        from PyPDF2 import PdfReader
    except Exception:
        return ""

    try:
        reader = PdfReader(str(file_path))
        parts: list[str] = []
        for page in reader.pages[:30]:
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t:
                parts.append(t)
            if sum(len(x) for x in parts) >= max_chars:
                break
        return preprocess_text("\n".join(parts), max_chars=max_chars)
    except Exception as e:
        logger.warning(f"pdf_text_extract_pypdf2_failed err={str(e)[:160]}")
        return ""


def _extract_page_range_text(file_path: str, start: int, end: int, max_chars: int) -> str:
    """Extract text from a specific page range of a PDF (worker for parallel extraction)."""
    try:
        import fitz
        _quiet_fitz()
    except Exception:
        return ""
    try:
        doc = fitz.open(str(file_path))
        parts: list[str] = []
        total = 0
        for i in range(start, min(end, len(doc))):
            try:
                t = doc[i].get_text("text") or ""
            except Exception:
                t = ""
            if t:
                parts.append(f"[Page {i+1}] {t}")
                total += len(t)
                if total >= max_chars:
                    break
        doc.close()
        return "\n".join(parts)
    except Exception:
        return ""


def extract_pdf_text_parallel(file_path: str, max_chars: int = 30000,
                               chunk_pages: int = 8, max_workers: int = 4) -> str:
    """Extract text from a large PDF by processing page ranges in parallel.

    Falls back to serial extraction for small files.
    """
    try:
        import fitz
        _quiet_fitz()
        doc = fitz.open(str(file_path))
        total_pages = len(doc)
        doc.close()
    except Exception:
        return extract_pdf_text_fast(file_path, max_chars)

    if total_pages <= chunk_pages:
        return extract_pdf_text_fast(file_path, max_chars)

    from concurrent.futures import ThreadPoolExecutor, as_completed

    # Split into page ranges
    chunks: list[tuple[int, int]] = []
    for start in range(0, total_pages, chunk_pages):
        end = min(start + chunk_pages, total_pages)
        chunks.append((start, end))

    # Per-chunk character budget
    per_chunk_chars = max(1000, max_chars // len(chunks))

    results: dict[int, str] = {}
    workers = min(max_workers, len(chunks))
    with ThreadPoolExecutor(max_workers=workers) as pool:
        futures = {}
        for order, (start, end) in enumerate(chunks):
            fut = pool.submit(_extract_page_range_text, file_path, start, end, per_chunk_chars)
            futures[fut] = order
        for fut in as_completed(futures):
            order = futures[fut]
            try:
                results[order] = fut.result()
            except Exception:
                results[order] = ""

    # Reassemble in page order
    text = "\n".join(results.get(i, "") for i in range(len(chunks)))
    return preprocess_text(text, max_chars=max_chars)


def extract_pdf_embedded_images(
    file_path: str,
    threshold: tuple[float, float] = (0.6, 0.6),
) -> str:
    """Extract embedded images from PDF pages and OCR them.

    Skips images whose width/page-width or height/page-height ratios
    fall below *threshold* (avoids OCR-ing small decorative images).
    """
    try:
        import fitz
        _quiet_fitz()
    except Exception:
        return ""

    try:
        doc = fitz.open(str(file_path))
        ocr_parts: list[str] = []
        for page in doc:
            page_w = page.rect.width
            page_h = page.rect.height
            for img_info in page.get_image_info(xrefs=True):
                w = img_info.get("width", 0)
                h = img_info.get("height", 0)
                if page_w > 0 and w / page_w < threshold[0]:
                    continue
                if page_h > 0 and h / page_h < threshold[1]:
                    continue
                try:
                    xref = img_info["xref"]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    # Handle rotation
                    if page.rotation != 0:
                        from PIL import Image

                        im = Image.open(BytesIO(image_bytes))
                        if page.rotation == 90:
                            im = im.transpose(Image.ROTATE_270)
                        elif page.rotation == 180:
                            im = im.transpose(Image.ROTATE_180)
                        elif page.rotation == 270:
                            im = im.transpose(Image.ROTATE_90)
                        buf = BytesIO()
                        im.save(buf, format="PNG")
                        image_bytes = buf.getvalue()
                    ocr_text = ocr_image(image_bytes)
                    if ocr_text:
                        ocr_parts.append(ocr_text)
                except Exception:
                    continue
        doc.close()
        return "\n".join(ocr_parts)
    except Exception as e:
        logger.warning(f"pdf_embedded_images_failed err={str(e)[:160]}")
        return ""


def render_pdf_pages(file_path: str, max_pages: int = 6) -> list[bytes]:
    """Render PDF pages as PNG images at 180 DPI (up to *max_pages* pages).

    Uses adaptive sampling for long documents: first 3 + samples from middle + last 2.
    """
    try:
        import fitz
        _quiet_fitz()
    except Exception:
        return []

    try:
        doc = fitz.open(str(file_path))
        total = doc.page_count
        if total == 0:
            doc.close()
            return []

        # Adaptive page selection for long documents
        if total <= max_pages:
            indices = list(range(total))
        elif total <= max_pages * 2:
            indices = list(range(max_pages))
        else:
            # First N pages + samples from middle + last N pages
            front = min(3, max_pages // 3)
            back = min(2, max_pages // 3)
            mid_count = max_pages - front - back
            if mid_count > 0 and total > front + back:
                step = max(1, (total - front - back) // mid_count)
                mid_indices = list(range(front, total - back, step))[:mid_count]
            else:
                mid_indices = []
            indices = list(range(front)) + mid_indices + list(range(total - back, total))
            indices = list(dict.fromkeys(indices))  # deduplicate while preserving order

        from concurrent.futures import ThreadPoolExecutor, as_completed

        def _render(i: int) -> tuple[int, bytes]:
            page = doc.load_page(i)
            return i, page.get_pixmap(dpi=180).tobytes("png")

        results: dict[int, bytes] = {}
        workers = min(len(indices), 4)
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {pool.submit(_render, i): i for i in indices}
            for fut in as_completed(futures):
                idx, data = fut.result()
                results[idx] = data
        doc.close()
        return [results[i] for i in sorted(results)]
    except Exception as e:
        logger.warning(f"pdf_render_failed err={str(e)[:160]}")
        return []


# ---------------------------------------------------------------------------
# PDF table extraction
# ---------------------------------------------------------------------------


def _clean_table_cell(value: Any) -> str:
    text = str(value or "").replace("\r", " ").replace("\n", " ").strip()
    text = re.sub(r"\s{2,}", " ", text)
    return text


def _table_to_markdown(rows: list[list[Any]]) -> str:
    cleaned_rows: list[list[str]] = []
    for row in rows:
        cells = [_clean_table_cell(cell) for cell in row]
        if any(cell for cell in cells):
            cleaned_rows.append(cells)
    if not cleaned_rows:
        return ""

    max_cols = max(len(row) for row in cleaned_rows)
    normalized = [row + [""] * (max_cols - len(row)) for row in cleaned_rows]
    header = normalized[0]
    divider = ["---"] * max_cols
    body = normalized[1:] or [[""] * max_cols]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(divider) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _normalize_table_rows(rows: list[list[Any]], max_rows: int = 30) -> tuple[list[list[str]], dict[str, Any]]:
    cleaned_rows: list[list[str]] = []
    raw_widths: list[int] = []
    for row in rows[:max_rows]:
        cells = [_clean_table_cell(cell) for cell in row]
        if any(cell for cell in cells):
            cleaned_rows.append(cells)
            raw_widths.append(len(cells))
    if not cleaned_rows:
        return [], {"column_count_inconsistent": False, "empty_header": False, "possible_truncation": False}

    max_cols = max(len(row) for row in cleaned_rows)
    normalized = [row + [""] * (max_cols - len(row)) for row in cleaned_rows]
    header = normalized[0] if normalized else []
    empty_header = bool(header) and not any(cell.strip() for cell in header)
    truncated_cells = sum(1 for row in normalized for cell in row if len(cell.strip()) == 1)
    checks = {
        "column_count_inconsistent": len(set(raw_widths)) > 1,
        "empty_header": empty_header,
        "possible_truncation": truncated_cells >= max(3, max_cols),
    }
    return normalized, checks


def _guess_table_caption(page: Any, bbox: Any) -> str:
    try:
        import fitz
        _quiet_fitz()
    except Exception:
        return ""
    if bbox is None:
        return ""
    try:
        rect = fitz.Rect(bbox)
        above = page.get_text("text", clip=fitz.Rect(rect.x0, max(0, rect.y0 - 90), rect.x1, rect.y0)).strip()
        below = page.get_text("text", clip=fitz.Rect(rect.x0, rect.y1, rect.x1, rect.y1 + 70)).strip()
    except Exception:
        return ""

    for candidate in reversed([line.strip() for line in above.splitlines() if line.strip()]):
        if 3 <= len(candidate) <= 80:
            return candidate
    for candidate in [line.strip() for line in below.splitlines() if line.strip()]:
        if 3 <= len(candidate) <= 80:
            return candidate
    return ""


def _table_summary_line(table: dict[str, Any]) -> str:
    headers = [str(x).strip() for x in table.get("headers", []) if str(x).strip()]
    dims = table.get("dimensions", {}) if isinstance(table.get("dimensions"), dict) else {}
    rows = int(dims.get("rows", len(table.get("rows", []) or [])) or 0)
    cols = int(dims.get("cols", len(headers) or 0) or 0)
    caption = str(table.get("caption") or "").strip()
    page = int(table.get("page", 0) or 0)
    bits = [f"第 {page} 页表格", f"{rows} 行 {cols} 列"]
    if headers:
        bits.append("表头：" + "/".join(headers[:6]))
    if caption:
        bits.append("标题：" + caption)
    return "（" + "，".join(bits) + "）"


def extract_pdf_tables_detailed(
    file_path: str,
    max_tables: int = 8,
    max_rows: int = 30,
) -> list[dict[str, Any]]:
    """Extract structured tables from PDF pages and validate basic structure."""
    try:
        import fitz
        _quiet_fitz()
    except Exception:
        return []

    tables_out: list[dict[str, Any]] = []
    doc = None
    try:
        doc = fitz.open(str(file_path))
        import contextlib
        _devnull = open(os.devnull, "w")
        try:
            for page_idx in range(len(doc)):
                page = doc[page_idx]
                try:
                    with contextlib.redirect_stderr(_devnull):
                        finder = page.find_tables()
                except Exception:
                    continue
                tables = list(getattr(finder, "tables", []) or [])
                for table_idx, table in enumerate(tables, start=1):
                    try:
                        rows = table.extract() or []
                    except Exception:
                        continue
                    normalized_rows, checks = _normalize_table_rows(rows, max_rows=max_rows)
                    if not normalized_rows:
                        continue

                    headers = normalized_rows[0]
                    body_rows = normalized_rows[1:]
                    if checks.get("empty_header"):
                        body_rows = normalized_rows
                        headers = []

                    item = {
                        "page": page_idx + 1,
                        "table_index": table_idx,
                        "caption": _guess_table_caption(page, getattr(table, "bbox", None)),
                        "headers": headers,
                        "rows": body_rows,
                        "dimensions": {
                            "rows": len(body_rows),
                            "cols": len(headers) or max(len(row) for row in normalized_rows),
                        },
                        "markdown": _table_to_markdown(normalized_rows),
                        "validation": checks,
                    }

                    if tables_out:
                        prev = tables_out[-1]
                        prev_dims = prev.get("dimensions", {}) if isinstance(prev.get("dimensions"), dict) else {}
                        same_width = int(prev_dims.get("cols", 0) or 0) == int(item["dimensions"]["cols"] or 0)
                        prev_page = int(prev.get("page", 0) or 0)
                        if (
                            checks.get("empty_header")
                            and same_width
                            and prev_page == page_idx
                        ):
                            prev_rows = prev.get("rows", [])
                            if isinstance(prev_rows, list):
                                prev_rows.extend(item["rows"])
                                prev["rows"] = prev_rows
                                prev["dimensions"]["rows"] = len(prev_rows)
                                prev["continued_pages"] = sorted(set([*prev.get("continued_pages", []), page_idx + 1]))
                                prev["validation"]["continued"] = True
                                prev["markdown"] = _table_to_markdown(
                                    ([prev.get("headers", [])] if prev.get("headers") else []) + prev_rows
                                )
                                continue

                    tables_out.append(item)
                    if len(tables_out) >= max_tables:
                        doc.close()
                        _devnull.close()
                        return tables_out
            doc.close()
        finally:
            _devnull.close()
    except Exception as e:
        logger.warning(f"pdf_table_extract_failed err={str(e)[:160]}")
        return []
    return tables_out


def parse_text_table_preview(
    text: str,
    max_rows: int = 8,
    max_cols: int = 12,
) -> list[dict[str, str]]:
    """Parse markdown/csv/hash-delimited tables from plain text."""
    raw_lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    lines = []
    for line in raw_lines:
        if line.lower().startswith(("标题:", "title:", "表格内容:", "caption:", "陈述:", "statement:", "标签:", "label:")):
            continue
        lines.append(line)

    table_lines = [line for line in lines if "|" in line or "#" in line or "\t" in line]
    if len(table_lines) < 2:
        return []

    delimiter = "|"
    sample = table_lines[0]
    if sample.count("|") >= 2:
        delimiter = "|"
    elif "#" in sample:
        delimiter = "#"
    elif "\t" in sample:
        delimiter = "\t"

    rows: list[list[str]] = []
    for line in table_lines:
        if delimiter == "|" and set(line.replace("|", "").replace("-", "").replace(" ", "")) == set():
            continue
        parts = [part.strip() for part in line.split(delimiter)]
        if delimiter == "|":
            parts = [part for part in parts if part != ""]
        if parts and not all(not p for p in parts):
            rows.append(parts[:max_cols])
        if len(rows) >= max_rows + 1:
            break
    if len(rows) < 2:
        return []

    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = rows[0]
    data_rows = rows[1 : 1 + max_rows]
    preview: list[dict[str, str]] = []
    for row in data_rows:
        item = {}
        for i, key in enumerate(header):
            col = key.strip() or f"col_{i+1}"
            item[col] = row[i].strip()
        preview.append(item)
    return preview


def extract_key_value_candidates(text: str, max_items: int = 10) -> list[dict[str, str]]:
    """Extract simple key-value pairs from OCR/plain text lines."""
    pairs: list[dict[str, str]] = []
    seen: set[tuple[str, str]] = set()
    for raw in str(text or "").splitlines():
        line = raw.strip().strip("|").strip()
        if len(line) < 3:
            continue
        if ":" in line:
            key, value = line.split(":", 1)
        elif "：" in line:
            key, value = line.split("：", 1)
        else:
            parts = re.split(r"\s{2,}", line)
            if len(parts) != 2:
                continue
            key, value = parts
        key = re.sub(r"\s+", " ", key).strip(" -")
        value = re.sub(r"\s+", " ", value).strip(" -")
        if len(key) < 2 or len(value) < 1:
            continue
        token = (key.lower(), value.lower())
        if token in seen:
            continue
        seen.add(token)
        pairs.append({"field": key, "value": value})
        if len(pairs) >= max_items:
            break
    return pairs


def extract_form_kv_by_text_sections(text: str, max_items: int = 10) -> list[dict[str, str]]:
    """Extract section-style form fields from OCR text using label lines and following content."""
    lines = [re.sub(r"\s+", " ", str(line or "")).strip() for line in str(text or "").splitlines()]
    lines = [line for line in lines if line]
    if not lines:
        return []

    label_re = re.compile(
        r"(?i)^([A-Z][A-Z0-9 .&()/\-]{2,}|[A-Za-z][A-Za-z0-9 .&()/\-]{2,})(?:[:：])$"
    )
    compact_label_re = re.compile(
        r"(?i)^(TO|FROM|DATE|PHONE|FAX|PAGES|REPORTINGPERIODS|TESTMARKETGEOGRAPHY|PRICEPOINT|SALESFORCEINVOLVEMENT|DISTRIBUTORS.*|CHAINS.*|INDEPENDENTS.*|ADVERTISING.*|TYPEOFPACKINGS|MANUFACTURER|BRAND)[:：]?$"
    )

    pairs: list[dict[str, str]] = []
    seen: set[tuple[str, str]] = set()
    idx = 0
    while idx < len(lines):
        line = lines[idx]
        is_label = bool(label_re.match(line) or compact_label_re.match(line.replace(" ", "")))
        if not is_label:
            idx += 1
            continue

        field = line.rstrip(":：").strip()
        value_lines: list[str] = []
        j = idx + 1
        while j < len(lines):
            nxt = lines[j]
            if label_re.match(nxt) or compact_label_re.match(nxt.replace(" ", "")):
                break
            value_lines.append(nxt)
            if len(" ".join(value_lines)) >= 260:
                break
            j += 1

        value = " ".join(value_lines).strip()
        if value:
            token = (field.lower(), value.lower())
            if token not in seen:
                seen.add(token)
                pairs.append({"field": field, "value": value})
        idx = max(idx + 1, j)
        if len(pairs) >= max_items:
            break

    return pairs


def extract_form_kv_by_reading_order(text: str, max_items: int = 12) -> list[dict[str, str]]:
    """Extract field/value pairs by OCR reading order with minimal normalization.

    This intentionally preserves raw-ish labels / values to better align with
    layout-annotation style benchmarks such as FUNSD.
    """
    lines = [re.sub(r"\s+", " ", str(line or "")).strip() for line in str(text or "").splitlines()]
    lines = [line for line in lines if line]
    if not lines:
        return []

    known_label_re = re.compile(
        r"(?i)^(TO|FROM|DATE|FAX NO\.?|FAX NUMBER|PHONE NUMBER|NUMBER OF PAGES INCLUDING COVER SHEET|SENDER/?PHONE NUMBER|SPECIAL INSTRUCTIONS|NOTE|MANUFACTURER|BRAND|TYPE OF PACKINGS|REPORTING PERIODS|TEST MARKET GEOGRAPHY|PRICE POINT|SALES FORCE INVOLVEMENT|DISTRIBUTORS.*|CHAINS.*|INDEPENDENTS.*|ADVERTISING.*)[:：]?$"
    )

    pairs: list[dict[str, str]] = []
    seen: set[tuple[str, str]] = set()

    def _add(field: str, value: str) -> None:
        f = re.sub(r"\s+", " ", str(field or "")).strip()
        v = re.sub(r"\s+", " ", str(value or "")).strip()
        if len(f) < 2 or len(v) < 1:
            return
        token = (f.lower(), v.lower())
        if token in seen:
            return
        seen.add(token)
        pairs.append({"field": f, "value": v})

    idx = 0
    while idx < len(lines):
        line = lines[idx]
        # Case 1: explicit inline label/value in same line.
        if ":" in line:
            field, value = line.split(":", 1)
            if value.strip():
                _add(field.strip() + ":", value.strip())
            else:
                j = idx + 1
                while j < len(lines) and not lines[j]:
                    j += 1
                if j < len(lines):
                    _add(field.strip() + ":", lines[j])
        # Case 2: label-only line, pair with next visible line.
        elif known_label_re.match(line):
            j = idx + 1
            while j < len(lines) and not lines[j]:
                j += 1
            if j < len(lines):
                _add(line, lines[j])

        if len(pairs) >= max_items:
            break
        idx += 1

    return pairs[:max_items]


def extract_form_kv_by_layout(
    ocr_items: list[dict[str, Any]],
    max_items: int = 12,
) -> list[dict[str, str]]:
    """Extract likely field-value pairs from OCR layout for scanned forms."""
    if not ocr_items:
        return []

    usable = [item for item in ocr_items if str(item.get("text") or "").strip()]
    if not usable:
        return []

    heights = [float(item.get("height") or 0.0) for item in usable if float(item.get("height") or 0.0) > 0]
    avg_height = sum(heights) / float(len(heights) or 1)
    line_gap = max(14.0, avg_height * 0.9)

    rows: list[list[dict[str, Any]]] = []
    for item in sorted(usable, key=lambda x: (float(x.get("cy") or 0.0), float(x.get("left") or 0.0))):
        cy = float(item.get("cy") or 0.0)
        placed = False
        for row in rows:
            row_cy = sum(float(x.get("cy") or 0.0) for x in row) / float(len(row) or 1)
            if abs(cy - row_cy) <= line_gap:
                row.append(item)
                placed = True
                break
        if not placed:
            rows.append([item])

    normalized_rows: list[list[dict[str, Any]]] = []
    for row in rows:
        normalized_rows.append(sorted(row, key=lambda x: float(x.get("left") or 0.0)))

    field_like = re.compile(
        r"(?i)^(to|from|date|fax|fax number|phone|phone number|number of pages|pages|sender|recipient|name|address|note|special instructions|company|country|market|product|subject|sample|brand|style|weights|adhesive|supplier|color|porosity|type|overall|dimension|positioning|request no\.?|completion|initiated by|country|product|received from|service request|licensee|operations|tar|nic|moist|preference|attributes)\b|.+[:：]$"
    )

    pairs: list[dict[str, str]] = []
    seen: set[tuple[str, str]] = set()

    def _add_pair(field: str, value: str) -> None:
        nonlocal pairs
        f = re.sub(r"\s+", " ", str(field or "")).strip(" -:：")
        v = re.sub(r"\s+", " ", str(value or "")).strip(" -:：")
        if len(f) < 2 or len(v) < 1:
            return
        token = (f.lower(), v.lower())
        if token in seen:
            return
        seen.add(token)
        pairs.append({"field": f, "value": v})

    def _join_text(items: list[dict[str, Any]]) -> str:
        parts = [str(x.get("text") or "").strip() for x in items if str(x.get("text") or "").strip()]
        return " ".join(parts).strip()

    def _looks_like_long_label(text: str) -> bool:
        s = re.sub(r"\s+", " ", str(text or "")).strip()
        if len(s) < 10:
            return False
        if field_like.search(s):
            return True
        return bool(
            re.search(
                r"(?i)(manufacturer|brand|type of packings|reporting periods|test market geography|price point|sales force involvement|acceptance|intro terms|intro deals|merchandising|effectiveness|distributors|chains|independents|advertising)",
                s,
            )
        )

    def _looks_like_value_row(text: str) -> bool:
        s = re.sub(r"\s+", " ", str(text or "")).strip()
        if len(s) < 2:
            return False
        if _looks_like_long_label(s):
            return False
        upper_ratio = sum(1 for ch in s if ch.isupper()) / float(sum(1 for ch in s if ch.isalpha()) or 1)
        return len(s) >= 6 and (upper_ratio < 0.75 or bool(re.search(r"\d|[.,()/$\"]", s)))

    for row_idx, row in enumerate(normalized_rows):
        texts = [str(x.get("text") or "").strip() for x in row]
        if not texts:
            continue

        # Pattern 1: inline "field: value" on the same OCR token.
        for text in texts:
            if ":" in text or "：" in text:
                sep = ":" if ":" in text else "："
                field, value = text.split(sep, 1)
                if value.strip():
                    _add_pair(field, value)

        # Pattern 2: left label, right value on the same row.
        for i, item in enumerate(row[:-1]):
            field = str(item.get("text") or "").strip()
            if not field:
                continue
            if field_like.search(field) or field.endswith(":") or field.endswith("："):
                right_group: list[dict[str, Any]] = []
                for candidate in row[i + 1 :]:
                    gap = float(candidate.get("left") or 0.0) - float(item.get("right") or 0.0)
                    if gap <= max(140.0, avg_height * 10):
                        right_group.append(candidate)
                value = _join_text(right_group)
                if value:
                    _add_pair(field, value)

        # Pattern 3: label on this row, value starts next row but aligned right.
        if row_idx + 1 < len(normalized_rows):
            next_row = normalized_rows[row_idx + 1]
            if next_row:
                left_item = row[0]
                field = str(left_item.get("text") or "").strip()
                next_text = _join_text(next_row[:4])
                if next_text and (field_like.search(field) or field.endswith(":") or field.endswith("：")):
                    next_left = float(next_row[0].get("left") or 0.0)
                    field_right = float(left_item.get("right") or 0.0)
                    if next_left >= field_right - 10:
                        _add_pair(field, next_text)

        # Pattern 4: left-column label row with value row directly underneath in same column block.
        if row_idx + 1 < len(normalized_rows):
            field_block = row
            next_row = normalized_rows[row_idx + 1]
            field_text = _join_text(field_block)
            if field_text and (field_like.search(field_text) or field_text.endswith(":") or field_text.endswith("：")):
                field_left = float(field_block[0].get("left") or 0.0)
                field_right = max(float(x.get("right") or 0.0) for x in field_block)
                candidate_items = [
                    x
                    for x in next_row
                    if float(x.get("left") or 0.0) >= field_left - avg_height
                    and float(x.get("left") or 0.0) <= field_right + max(160.0, avg_height * 12)
                ]
                candidate_text = _join_text(candidate_items)
                if candidate_text:
                    _add_pair(field_text, candidate_text)

        # Pattern 5: long label row followed by one or more paragraph-style value rows.
        field_text = _join_text(row)
        if field_text and _looks_like_long_label(field_text) and row_idx + 1 < len(normalized_rows):
            value_lines: list[str] = []
            field_left = min(float(x.get("left") or 0.0) for x in row)
            field_right = max(float(x.get("right") or 0.0) for x in row)
            max_follow_rows = min(len(normalized_rows), row_idx + 4)
            for next_idx in range(row_idx + 1, max_follow_rows):
                next_row = normalized_rows[next_idx]
                next_text = _join_text(next_row)
                if not next_text:
                    continue
                if _looks_like_long_label(next_text):
                    break
                next_left = min(float(x.get("left") or 0.0) for x in next_row)
                next_right = max(float(x.get("right") or 0.0) for x in next_row)
                if next_left > field_right + max(120.0, avg_height * 8):
                    break
                if next_right < field_left - max(40.0, avg_height * 2):
                    continue
                if _looks_like_value_row(next_text):
                    value_lines.append(next_text)
                if len(" ".join(value_lines)) >= 360:
                    break
            if value_lines:
                _add_pair(field_text, " ".join(value_lines))

        if len(pairs) >= max_items:
            break

    return pairs[:max_items]


def extract_pdf_tables_structured(
    file_path: str,
    max_tables: int = 8,
    max_rows: int = 30,
    max_chars: int = 8000,
) -> str:
    """Return a compact summary of detected PDF tables for prompt injection."""
    tables = extract_pdf_tables_detailed(file_path, max_tables=max_tables, max_rows=max_rows)
    if not tables:
        return ""
    chunks: list[str] = []
    total_chars = 0
    for table in tables:
        chunk = _table_summary_line(table)
        chunks.append(chunk)
        total_chars += len(chunk)
        if total_chars >= max_chars:
            break
    return preprocess_text("\n".join(chunks), max_chars=max_chars)


# ---------------------------------------------------------------------------
# DOCX structured extraction
# ---------------------------------------------------------------------------


def extract_docx_structured(file_path: str, max_chars: int = 20000) -> str:
    """Extract text, tables, and embedded images from a .docx file.

    Uses python-docx to walk paragraphs and tables in document order,
    and OCRs embedded images found in paragraphs.
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except Exception:
        return ""

    try:
        doc = Document(str(file_path))
        parts: list[str] = []

        for block in doc.element.body:
            tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

            if tag == "p":
                # Paragraph
                para = _paragraph_for_element(doc, block)
                if para is None:
                    continue
                text = para.text.strip()
                if text:
                    try:
                        style_name = str(getattr(getattr(para, "style", None), "name", "") or "").lower()
                    except Exception:
                        style_name = ""
                    if style_name.startswith("heading"):
                        level_match = re.search(r"(\d+)", style_name)
                        level = int(level_match.group(1)) if level_match else 1
                        level = max(1, min(level, 6))
                        parts.append("#" * level + " " + text)
                    else:
                        parts.append(text)
                # Check for embedded images in this paragraph
                images_text = _extract_images_from_paragraph(para, doc)
                if images_text:
                    parts.append(images_text)

            elif tag == "tbl":
                # Table
                table_text = _extract_table_text(block)
                if table_text:
                    parts.append(table_text)

        result = "\n\n".join(parts)
        return preprocess_text(result, max_chars=max_chars)
    except Exception as e:
        logger.warning(f"docx_structured_failed err={str(e)[:160]}")
        return ""


def _paragraph_for_element(doc, element) -> Any:
    """Find the Paragraph object matching an XML element."""
    from docx.text.paragraph import Paragraph

    for para in doc.paragraphs:
        if para._element is element:
            return para
    return None


def _extract_images_from_paragraph(para, doc) -> str:
    """Extract and OCR embedded images in a paragraph."""
    try:
        from docx.oxml.ns import qn
    except Exception:
        return ""

    pics = para._element.xpath(".//*[local-name()='pic']")
    if not pics:
        return ""

    # Build rId → image part lookup
    image_parts: dict[str, Any] = {}
    for rel in doc.part.rels.values():
        if "image" in (rel.reltype or ""):
            image_parts[rel.rId] = rel.target_part

    texts: list[str] = []
    for pic in pics:
        blips = pic.xpath(".//*[local-name()='blip']")
        for blip in blips:
            embed = blip.get(qn("r:embed"))
            if not embed or embed not in image_parts:
                continue
            try:
                image_bytes = image_parts[embed].blob
                ocr_text = ocr_image(image_bytes)
                if ocr_text:
                    texts.append(ocr_text)
            except Exception:
                continue
    return "\n".join(texts) if texts else ""


def _extract_table_text(tbl_element) -> str:
    """Extract text from a table XML element."""
    try:
        from docx.oxml.ns import qn
    except Exception:
        return ""

    rows: list[str] = []
    for row in tbl_element.findall(qn("w:tr")):
        cells: list[str] = []
        for cell in row.findall(qn("w:tc")):
            cell_texts: list[str] = []
            for p in cell.findall(qn("w:p")):
                t = "".join(t_node.text or "" for t_node in p.findall(".//" + qn("w:t")))
                if t.strip():
                    cell_texts.append(t.strip())
            cells.append(" ".join(cell_texts))
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows) if rows else ""


# ---------------------------------------------------------------------------
# PPTX structured extraction
# ---------------------------------------------------------------------------


def extract_pptx_structured(file_path: str, max_chars: int = 20000) -> str:
    """Extract text and embedded images from a .pptx file.

    Uses python-pptx to walk shapes on each slide (sorted positionally),
    and OCRs embedded images found in picture shapes or groups.
    """
    try:
        from pptx import Presentation
        from pptx.shapes.picture import Picture
        from pptx.shapes.group import GroupShapes
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:
        return ""

    try:
        prs = Presentation(str(file_path))
        parts: list[str] = []

        for slide in prs.slides:
            slide_parts: list[str] = []
            title = ""
            try:
                title = str(slide.shapes.title.text or "").strip() if slide.shapes.title else ""
            except Exception:
                title = ""
            if title:
                slide_parts.append("# " + title)

            # Sort shapes top-to-bottom, left-to-right
            shapes_sorted = sorted(
                list(slide.shapes),
                key=lambda s: (getattr(s, "top", 0) or 0, getattr(s, "left", 0) or 0),
            )

            for shape in shapes_sorted:
                shape_text = _extract_pptx_shape_text(shape)
                if shape_text:
                    if title and shape_text.strip() == title:
                        continue
                    slide_parts.append(shape_text)

            if slide_parts:
                parts.append("\n".join(slide_parts))

        result = "\n\n".join(parts)
        return preprocess_text(result, max_chars=max_chars)
    except Exception as e:
        logger.warning(f"pptx_structured_failed err={str(e)[:160]}")
        return ""


def _extract_pptx_shape_text(shape) -> str:
    """Extract text from a single PPTX shape (text, table, picture, group)."""
    from pptx.shapes.picture import Picture
    from pptx.shapes.group import GroupShapes
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    results: list[str] = []

    # Text frame
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                results.append(t)

    # Table
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            results.append(" | ".join(cells))

    # Embedded image (picture)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image_bytes = shape.image.blob
            ocr_text = ocr_image(image_bytes)
            if ocr_text:
                results.append(ocr_text)
        except Exception:
            pass

    # Group shape — recurse
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                child_text = _extract_pptx_shape_text(child)
                if child_text:
                    results.append(child_text)
        except Exception:
            pass

    return "\n".join(results) if results else ""


# ---------------------------------------------------------------------------
# Image table structure extraction
# ---------------------------------------------------------------------------


def extract_table_structure_from_image(
    ocr_items: list[dict[str, Any]],
    image_bytes: bytes | None = None,
) -> list[dict[str, Any]]:
    """Extract structured table data from an image.

    Uses OCR spatial analysis to determine row / column boundaries, then
    optionally crops each cell and runs OCR on the individual cell image for
    higher-accuracy text recovery.

    Returns a list of table dicts with the same schema as
    :func:`extract_pdf_tables_detailed`.
    """
    if not ocr_items or len(ocr_items) < 4:
        return []

    # ── 1.  Filter items with valid text and known centres ──────────────────
    items: list[dict[str, Any]] = []
    for it in ocr_items:
        text = str(it.get("text") or "").strip()
        cx = it.get("cx")
        cy = it.get("cy")
        if text and cx is not None and cy is not None:
            items.append({"text": text, "cx": float(cx), "cy": float(cy)})

    if len(items) < 4:
        return []

    # ── 2.  Cluster items into rows (y-coordinate proximity) ─────────────────
    rows = _cluster_ocr_items_into_rows(items)
    if len(rows) < 2:
        return []

    # ── 3.  Determine column boundaries from x-coordinate clusters ───────────
    col_centers, col_bounds = _compute_col_boundaries(items)
    n_cols = len(col_centers)
    if n_cols < 2:
        return []

    # ── 4.  Try grid-detection when image has enough resolution ──────────────
    table_rows: list[list[str]] = []
    if image_bytes is not None:
        from PIL import Image as PILImage
        _img_w, _img_h = PILImage.open(BytesIO(image_bytes)).size
        if min(_img_w, _img_h) >= 200:
            grid_rows = _extract_via_grid_detection(image_bytes, items)
            if grid_rows and len(grid_rows) >= 2:
                table_rows = grid_rows

    # ── 5.  Fall back to OCR spatial assignment ──────────────────────────────
    if not table_rows:
        table_rows = _assign_ocr_to_cells(rows, col_bounds, n_cols)

    if len(table_rows) < 2:
        return []

    # ── 6.  Normalise column count, split header/body, build Markdown ────────
    headers = list(table_rows[0])
    body = table_rows[1:]

    max_cols = max(len(headers), max((len(r) for r in body), default=0))
    if max_cols == 0:
        return []

    headers = headers + [""] * (max_cols - len(headers))
    body = [r + [""] * (max_cols - len(r)) for r in body]

    md_lines = ["| " + " | ".join(headers) + " |"]
    md_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in body:
        md_lines.append("| " + " | ".join(row) + " |")

    return [
        {
            "page": 0,
            "table_index": 1,
            "caption": "",
            "headers": headers,
            "rows": body,
            "dimensions": {"rows": len(body), "cols": max_cols},
            "markdown": "\n".join(md_lines),
            "validation": {"method": "ocr_cell_level" if image_bytes else "ocr_spatial"},
        }
    ]


def _extract_via_grid_detection(
    image_bytes: bytes,
    ocr_items: list[dict[str, Any]],
) -> list[list[str]]:
    """Extract table cells using OpenCV grid-line detection + OCR mapping.

    Detects horizontal and vertical grid lines via morphological operations,
    then maps existing OCR items into cells.  Only falls back to per-cell
    OCR for cells that received no text from the global OCR pass.

    Works best on clean rendered tables with visible borders (PubTabNet).
    Returns ``[]`` when fewer than 2 rows or 2 columns are found.
    """
    import cv2
    import numpy as np
    from PIL import Image as PILImage

    img = PILImage.open(BytesIO(image_bytes)).convert("L")
    img_np = np.array(img)
    h, w = img_np.shape

    # ── 1.  Binarise ─────────────────────────────────────────────────────
    _, binary = cv2.threshold(img_np, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

    # ── 2.  Detect horizontal & vertical lines ────────────────────────────
    h_kernel_len = max(4, min(30, int(w * 0.03)))
    v_kernel_len = max(3, min(20, int(h * 0.05)))

    h_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (h_kernel_len, 1))
    v_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, v_kernel_len))

    h_lines = cv2.morphologyEx(binary, cv2.MORPH_OPEN, h_kernel)
    v_lines = cv2.morphologyEx(binary, cv2.MORPH_OPEN, v_kernel)

    # ── 3.  Find row / column edges from projections ──────────────────────
    # Filter vertical lines by connected-component height — real table
    # borders span most of the table, text strokes are short.
    v_dilated = cv2.dilate(v_lines, np.ones((1, 3), np.uint8), iterations=1)
    v_labels = cv2.connectedComponents(v_dilated, connectivity=8)[1]
    v_filtered = np.zeros_like(v_lines)
    for label in range(1, int(v_labels.max()) + 1):
        mask = (v_labels == label).astype(np.uint8)
        ys = np.where(mask.any(axis=1))[0]
        comp_height = float(ys[-1] - ys[0] + 1) if len(ys) > 0 else 0.0
        if comp_height >= h * 0.15:
            v_filtered[mask > 0] = 255

    h_dilated = cv2.dilate(h_lines, np.ones((3, 1), np.uint8), iterations=1)
    h_labels = cv2.connectedComponents(h_dilated, connectivity=8)[1]
    h_filtered = np.zeros_like(h_lines)
    for label in range(1, int(h_labels.max()) + 1):
        mask = (h_labels == label).astype(np.uint8)
        xs = np.where(mask.any(axis=0))[0]
        comp_width = float(xs[-1] - xs[0] + 1) if len(xs) > 0 else 0.0
        if comp_width >= w * 0.30:
            h_filtered[mask > 0] = 255

    h_proj = np.sum(h_filtered, axis=1, dtype=np.float64) / 255.0
    v_proj = np.sum(v_filtered, axis=0, dtype=np.float64) / 255.0

    h_thresh = max(float(np.max(h_proj)) * 0.25, 2.0)
    v_thresh = max(float(np.max(v_proj)) * 0.25, 2.0)

    row_positions = _find_line_positions(h_proj, h_thresh)
    col_positions = _find_line_positions(v_proj, v_thresh)

    if len(row_positions) < 2 or len(col_positions) < 2:
        return []

    row_edges = [0] + sorted(row_positions) + [h]
    col_edges = [0] + sorted(col_positions) + [w]

    n_rows = len(row_edges) - 1
    n_cols = len(col_edges) - 1

    # ── 4.  Map OCR items into grid cells (fast path) ─────────────────────
    cell_texts: list[list[list[str]]] = [[[] for _ in range(n_cols)] for _ in range(n_rows)]
    for it in ocr_items:
        cx = float(it.get("cx", 0))
        cy = float(it.get("cy", 0))
        text = str(it.get("text") or "").strip()
        if not text or cx <= 0 or cy <= 0:
            continue
        ri = _find_row_index(cy, row_edges)
        ci = _find_column_index(cx, col_edges)
        if ri is not None and ci is not None:
            cell_texts[ri][ci].append(text)

    # ── 5.  Assemble — cell-crop OCR only for empty cells ─────────────────
    table_rows: list[list[str]] = []
    for r in range(n_rows):
        row_cells: list[str] = []
        for c in range(n_cols):
            texts = cell_texts[r][c]
            if texts:
                cell = " ".join(texts).strip()
            else:
                x1, y1 = int(col_edges[c]), int(row_edges[r])
                x2, y2 = int(col_edges[c + 1]), int(row_edges[r + 1])
                cw, ch = x2 - x1, y2 - y1
                if min(cw, ch) < 8:
                    cell = ""
                else:
                    cell_crop = img_np[y1:y2, x1:x2]
                    if min(cw, ch) < 40:
                        cell_crop = cv2.resize(
                            cell_crop, (cw * 3, ch * 3),
                            interpolation=cv2.INTER_CUBIC,
                        )
                    buf = BytesIO()
                    PILImage.fromarray(cell_crop).save(buf, format="PNG")
                    cell = ocr_image(buf.getvalue()).strip()
                    cell = re.sub(r"\s+", " ", cell).strip()
            row_cells.append(cell)
        if any(c for c in row_cells):
            table_rows.append(row_cells)

    return table_rows


def _find_line_positions(
    proj: "np.ndarray",
    threshold: float,
    min_gap: int = 6,
) -> list[int]:
    """Extract line positions from a 1-D projection above *threshold*."""
    import numpy as np

    above = [int(i) for i, v in enumerate(proj) if v >= threshold]
    if not above:
        return []

    clusters: list[list[int]] = []
    current = [above[0]]
    for pos in above[1:]:
        if pos - current[-1] <= min_gap:
            current.append(pos)
        else:
            clusters.append(current)
            current = [pos]
    clusters.append(current)

    return [int(np.median(np.array(cluster, dtype=np.float64))) for cluster in clusters]


def _find_row_index(cy: float, row_edges: list[int]) -> int | None:
    """Return the row index whose y-interval contains *cy*."""
    for i in range(len(row_edges) - 1):
        if row_edges[i] <= int(cy) <= row_edges[i + 1]:
            return i
    return None


def _find_column_index(cx: float, col_edges: list[int]) -> int | None:
    """Return the column index whose x-interval contains *cx*."""
    for i in range(len(col_edges) - 1):
        if col_edges[i] <= int(cx) <= col_edges[i + 1]:
            return i
    return None


def _assign_ocr_to_cells(
    rows: list[list[dict[str, Any]]],
    col_bounds: list[float],
    n_cols: int,
) -> list[list[str]]:
    """Fallback: assign OCR items to grid cells by position."""
    table_rows: list[list[str]] = []
    for row_items in rows:
        cells = ["" for _ in range(n_cols)]
        for it in row_items:
            ci = _find_column(it["cx"], col_bounds)
            if ci is not None:
                prev = cells[ci]
                cells[ci] = (prev + " " + it["text"]).strip() if prev else it["text"]
        if any(c.strip() for c in cells):
            table_rows.append(cells)
    return table_rows


def _cluster_ocr_items_into_rows(
    items: list[dict[str, Any]],
) -> list[list[dict[str, Any]]]:
    """Cluster OCR items into visual rows by y-coordinate proximity."""
    if not items:
        return []

    sorted_items = sorted(items, key=lambda x: x["cy"])

    # Compute adaptive row-gap threshold from median adjacent gap
    cy_list = [it["cy"] for it in sorted_items]
    gaps = [cy_list[i + 1] - cy_list[i] for i in range(len(cy_list) - 1)]
    gaps = [g for g in gaps if g > 1.0]
    if gaps:
        median_gap = sorted(gaps)[len(gaps) // 2]
        threshold = max(6.0, median_gap * 0.70)
    else:
        threshold = 10.0

    rows: list[list[dict[str, Any]]] = []
    for item in sorted_items:
        placed = False
        for row in rows:
            row_cy = sum(r["cy"] for r in row) / len(row)
            if abs(item["cy"] - row_cy) <= threshold:
                row.append(item)
                placed = True
                break
        if not placed:
            rows.append([item])

    # Sort left-to-right within each row
    for row in rows:
        row.sort(key=lambda x: x["cx"])

    # ── Merge sparse header rows ──────────────────────────────────────────
    # A single-item row right above a multi-item row is often a colspan
    # header (e.g. "Age (years)" spanning 2 columns) that OCR placed on
    # a slightly different baseline.
    merged_rows = _merge_sparse_rows(rows, threshold)

    return merged_rows


def _merge_sparse_rows(
    rows: list[list[dict[str, Any]]],
    row_threshold: float,
) -> list[list[dict[str, Any]]]:
    """Merge adjacent rows where one is very sparse compared to its neighbour."""
    if len(rows) < 2:
        return rows

    result: list[list[dict[str, Any]]] = [list(rows[0])]
    for idx in range(1, len(rows)):
        prev = result[-1]
        curr = rows[idx]
        prev_cy = sum(r["cy"] for r in prev) / len(prev)
        curr_cy = sum(r["cy"] for r in curr) / len(curr)
        gap = curr_cy - prev_cy

        # Merge if the gap is less than 2× the intra-row threshold AND
        # one of the two rows has ≤2 items while the other has ≥3.
        should_merge = (
            gap <= row_threshold * 2.5
            and (len(prev) <= 2) != (len(curr) <= 2)
            and max(len(prev), len(curr)) >= 3
        )
        if should_merge:
            prev.extend(curr)
            prev.sort(key=lambda x: x["cx"])
        else:
            result.append(list(curr))

    return result


def _compute_col_boundaries(
    items: list[dict[str, Any]],
) -> tuple[list[float], list[float]]:
    """Compute column centres and boundaries from x-coordinate clusters."""
    cx_list = sorted(set(it["cx"] for it in items))
    if len(cx_list) <= 1:
        return cx_list, [float("-inf"), float("inf")]

    gaps = [cx_list[i + 1] - cx_list[i] for i in range(len(cx_list) - 1)]
    if not gaps:
        return cx_list, [float("-inf"), float("inf")]

    # Use 75th-percentile gap as threshold — more robust than median
    p75_idx = int(len(gaps) * 0.75)
    threshold = max(10.0, sorted(gaps)[p75_idx] * 0.45) if gaps else 10.0

    # DBSCAN-like single-pass clustering
    clusters: list[list[float]] = []
    for cx in cx_list:
        placed = False
        for cluster in clusters:
            if abs(cx - sum(cluster) / len(cluster)) <= threshold:
                cluster.append(cx)
                placed = True
                break
        if not placed:
            clusters.append([cx])

    col_centers = sorted(sum(c) / len(c) for c in clusters)
    col_bounds = [float("-inf")]
    for i in range(len(col_centers) - 1):
        col_bounds.append((col_centers[i] + col_centers[i + 1]) / 2.0)
    col_bounds.append(float("inf"))

    return col_centers, col_bounds


def _find_column(cx: float, col_bounds: list[float]) -> Optional[int]:
    """Return the column index whose interval contains *cx*, or ``None``."""
    for c in range(len(col_bounds) - 1):
        if col_bounds[c] <= cx < col_bounds[c + 1]:
            return c
    return None


# ---------------------------------------------------------------------------
# Quality metrics
# ---------------------------------------------------------------------------


def compute_quality_metrics(result: Any) -> dict:
    fields = ["title", "keywords", "summary", "key_points", "data", "conclusion", "prediction"]
    filled = 0
    for f in fields:
        v = result.get(f) if isinstance(result, dict) else None
        ok = False
        if isinstance(v, str):
            ok = bool(v.strip())
        elif isinstance(v, list) or isinstance(v, dict):
            ok = bool(v)
        if ok:
            filled += 1
    required = ["title", "summary", "key_points"]
    missing_required = []
    for key in required:
        value = result.get(key) if isinstance(result, dict) else None
        if isinstance(value, str) and value.strip():
            continue
        if isinstance(value, list) and value:
            continue
        missing_required.append(key)
    return {
        "filled_fields": filled,
        "field_ratio": filled / float(len(fields) or 1),
        "missing_required": missing_required,
        "required_ok": not missing_required,
    }


def choose_better_extraction(a: dict, b: dict) -> dict:
    """Choose the stronger extraction result using required fields first."""
    qa = compute_quality_metrics(a)
    qb = compute_quality_metrics(b)

    score_a = (
        (3 if qa.get("required_ok") else 0)
        + int(qa.get("filled_fields", 0))
        + min(len(str(a.get("summary") or "").strip()) / 200.0, 1.0)
    )
    score_b = (
        (3 if qb.get("required_ok") else 0)
        + int(qb.get("filled_fields", 0))
        + min(len(str(b.get("summary") or "").strip()) / 200.0, 1.0)
    )
    return b if score_b >= score_a else a


def _similarity_tokens(text: str) -> set[str]:
    parts = re.findall(r"[a-z0-9]{2,}|[\u4e00-\u9fff]", (text or "").lower())
    return {p for p in parts if p.strip()}


def _text_similarity(a: str, b: str) -> float:
    ta = _similarity_tokens(a)
    tb = _similarity_tokens(b)
    jaccard = 0.0
    if ta or tb:
        jaccard = len(ta & tb) / float(len(ta | tb) or 1)
    seq = SequenceMatcher(None, (a or "").strip(), (b or "").strip()).ratio()
    return max(jaccard, seq)


def deduplicate_analysis_results(results: list[dict], threshold: float = 0.82) -> list[dict]:
    """Merge near-duplicate extraction results to reduce repeated sources downstream."""
    deduped: list[dict] = []
    for item in results or []:
        if not isinstance(item, dict):
            continue
        candidate = dict(item)
        file_name = str(candidate.get("_file") or "").strip()
        title = str(candidate.get("title") or "").strip()
        summary = str(candidate.get("summary") or "").strip()
        key_points = candidate.get("key_points") if isinstance(candidate.get("key_points"), list) else []
        basis = " ".join([title, summary, " ".join(str(x) for x in key_points[:5])]).strip()
        if not basis:
            deduped.append(candidate)
            continue

        merged = False
        for idx, existing in enumerate(deduped):
            e_title = str(existing.get("title") or "").strip()
            e_summary = str(existing.get("summary") or "").strip()
            e_key_points = existing.get("key_points") if isinstance(existing.get("key_points"), list) else []
            existing_basis = " ".join([e_title, e_summary, " ".join(str(x) for x in e_key_points[:5])]).strip()
            if not existing_basis:
                continue

            sim = _text_similarity(basis, existing_basis)
            if sim < threshold:
                continue

            better = choose_better_extraction(existing, candidate)
            weaker = candidate if better is existing else existing
            merged_item = dict(better)

            merged_files = []
            existing_merged = existing.get("_merged_files", []) if isinstance(existing.get("_merged_files"), list) else []
            candidate_merged = candidate.get("_merged_files", []) if isinstance(candidate.get("_merged_files"), list) else []
            for name in [*existing_merged, *candidate_merged]:
                if isinstance(name, str) and name.strip() and name not in merged_files:
                    merged_files.append(name)
            for name in [str(existing.get("_file") or "").strip(), str(candidate.get("_file") or "").strip()]:
                if name and name not in merged_files:
                    merged_files.append(name)
            merged_item["_merged_files"] = merged_files
            merged_item["_dedup_score"] = round(sim, 4)

            if not merged_item.get("summary") and weaker.get("summary"):
                merged_item["summary"] = weaker.get("summary")
            if not merged_item.get("title") and weaker.get("title"):
                merged_item["title"] = weaker.get("title")
            if not merged_item.get("data") and weaker.get("data"):
                merged_item["data"] = weaker.get("data")

            merged_points: list[str] = []
            for source in (
                existing.get("key_points") if isinstance(existing.get("key_points"), list) else [],
                candidate.get("key_points") if isinstance(candidate.get("key_points"), list) else [],
            ):
                for point in source:
                    text = str(point or "").strip()
                    if text and text not in merged_points:
                        merged_points.append(text)
            if merged_points:
                merged_item["key_points"] = merged_points[:7]

            deduped[idx] = merged_item
            merged = True
            break

        if not merged:
            if file_name:
                candidate["_merged_files"] = [file_name]
            deduped.append(candidate)
    return deduped
